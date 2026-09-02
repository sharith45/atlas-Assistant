import os
import sys
import json
import argparse
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Configuración de codificación UTF-8 en consola de Windows
if sys.platform == "win32":
    try:
        sys.stdout.reconfigure(encoding="utf-8")
        sys.stderr.reconfigure(encoding="utf-8")
    except Exception:
        pass

def hex_to_rgb(hex_str):
    """Convierte un color en hexadecimal (con o sin #) a RGBColor de python-pptx."""
    if not hex_str or not isinstance(hex_str, str):
        return RGBColor(0x1E, 0x3A, 0x8A)
    hex_str = hex_str.lstrip('#').strip()
    if len(hex_str) > 6:
        hex_str = hex_str[:6]
    elif len(hex_str) == 3:
        hex_str = ''.join([c*2 for c in hex_str])
    elif len(hex_str) < 6:
        hex_str = hex_str.ljust(6, '0')
    try:
        r = int(hex_str[0:2], 16)
        g = int(hex_str[2:4], 16)
        b = int(hex_str[4:6], 16)
        return RGBColor(r, g, b)
    except Exception:
        return RGBColor(0x1E, 0x3A, 0x8A)

# Registro de Temas Visuales Predefinidos (Estilos Canva / Modernos)
TEMAS = {
    "dark_bento": {
        "bg": "0F172A",
        "card_bg": "1E293B",
        "card_border": "334155",
        "text_primary": "F8FAFC",
        "text_muted": "94A3B8",
        "accent_1": "38BDF8", # Cyan
        "accent_2": "C084FC", # Purple
        "accent_3": "34D399", # Emerald
        "font_family": "Calibri"
    },
    "canva_creative": {
        "bg": "F8FAFC",
        "card_bg": "EEF2FF",
        "card_border": "C7D2FE",
        "text_primary": "0F172A",
        "text_muted": "475569",
        "accent_1": "4F46E5", # Indigo
        "accent_2": "059669", # Emerald
        "accent_3": "E11D48", # Rose
        "font_family": "Calibri"
    },
    "corporate_executive": {
        "bg": "FFFFFF",
        "card_bg": "F8FAFC",
        "card_border": "E2E8F0",
        "text_primary": "0F172A",
        "text_muted": "64748B",
        "accent_1": "1E3A8A", # Navy Blue
        "accent_2": "D97706", # Amber/Gold
        "accent_3": "2563EB", # Sapphire
        "font_family": "Calibri"
    },
    "neon_cyberpunk": {
        "bg": "0B0F19",
        "card_bg": "131A29",
        "card_border": "EC4899", # Magenta Neon
        "text_primary": "FFFFFF",
        "text_muted": "9CA3AF",
        "accent_1": "06B6D4", # Cyan Neon
        "accent_2": "EC4899", # Pink Neon
        "accent_3": "FACC15", # Yellow Neon
        "font_family": "Calibri"
    },
    "minimal_editorial": {
        "bg": "FAF8F5",
        "card_bg": "FFFFFF",
        "card_border": "E5E0D8",
        "text_primary": "1C1917",
        "text_muted": "78716C",
        "accent_1": "B45309", # Warm Gold
        "accent_2": "44403C", # Stone
        "accent_3": "D97706", # Amber
        "font_family": "Georgia"
    },
    "vibrant_roadmap": {
        "bg": "F1F5F9",
        "card_bg": "8B5CF6",
        "card_border": "CBD5E1",
        "text_primary": "0F172A",
        "text_muted": "475569",
        "accent_1": "8B5CF6", # Violet
        "accent_2": "EC4899", # Pink
        "accent_3": "3B82F6", # Blue
        "font_family": "Calibri"
    }
}

from pptx.oxml import parse_xml

# XMLs de Transiciones Nativas de PowerPoint
TRANSICIONES_XML = {
    "fade": '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:fade/></p:transition>',
    "push": '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:push dir="r"/></p:transition>',
    "wipe": '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:wipe dir="l"/></p:transition>',
    "zoom": '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:zoom dir="in"/></p:transition>',
    "cover": '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:cover dir="r"/></p:transition>',
    "dissolve": '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:dissolve/></p:transition>',
    "wheel": '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:wheel spoke="4"/></p:transition>',
    "split": '<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" spd="med"><p:split dir="in" orient="horz"/></p:transition>'
}

def apply_transition(slide, trans_type="fade"):
    if not trans_type or trans_type == "none":
        return
    trans_type = trans_type.lower()
    xml_str = TRANSICIONES_XML.get(trans_type, TRANSICIONES_XML["fade"])
    try:
        transition_elm = parse_xml(xml_str)
        slide._element.append(transition_elm)
    except Exception:
        pass

def add_solid_bg(slide, prs, color_hex, slide_idx=0, accent_hex=None, accent_sec_hex=None):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = hex_to_rgb(color_hex)
    bg.line.fill.background()

    # Fondo dinámico y variado por índice de diapositiva (Sin repetición monótona)
    if slide_idx == 0:
        # Portada: Gran arco superior de acento
        if accent_hex:
            circ1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10.2), Inches(-2.0), Inches(5.0), Inches(5.0))
            circ1.fill.solid()
            circ1.fill.fore_color.rgb = hex_to_rgb(accent_hex)
            circ1.line.fill.background()
    elif slide_idx % 3 == 1:
        # Trama de puntos halftone en esquina superior derecha
        add_halftone_pattern(slide, start_x=11.2, start_y=0.4, cols=5, rows=4, dot_color=accent_hex or "94A3B8")
    elif slide_idx % 3 == 2:
        # Trama de puntos en esquina inferior izquierda
        add_halftone_pattern(slide, start_x=0.5, start_y=5.2, cols=5, rows=4, dot_color=accent_sec_hex or "38BDF8")

    # Franja superior discreta de acento
    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.06))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = hex_to_rgb(accent_hex or "38BDF8")
    top_bar.line.fill.background()
    return bg

def add_drop_shadow(shape, offset_x=2, offset_y=2, blur=4, opacity=0.15, color_hex="000000"):
    """Agrega una sombra drop shadow profesional a una forma mediante manipulación OpenXML."""
    try:
        color_clean = color_hex.lstrip('#').strip()
        if len(color_clean) != 6:
            color_clean = "000000"
        dist_val = int((offset_x + offset_y) * 12700 / 2)
        blur_val = int(blur * 12700)
        alpha_val = int(opacity * 100000)

        shadow_xml = f'''
        <a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">
            <a:outerShdw blurRad="{blur_val}" dist="{dist_val}" dir="2700000" algn="tl" rotWithShape="0">
                <a:srgbClr val="{color_clean}">
                    <a:alpha val="{alpha_val}"/>
                </a:srgbClr>
            </a:outerShdw>
        </a:effectLst>
        '''
        shape.element.spPr.append(parse_xml(shadow_xml))
    except Exception:
        pass

def add_footer(slide, prs, current_page, total_pages, title_text="", font_name="Calibri", text_hex="94A3B8"):
    """Añade pie de página elegante con paginación y título."""
    if current_page <= 1:
        return
    txBox = slide.shapes.add_textbox(Inches(1.0), Inches(7.0), Inches(11.333), Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_top = tf.margin_bottom = tf.margin_left = tf.margin_right = 0
    p = tf.paragraphs[0]
    p.text = f"{title_text}   •   Diapositiva {current_page} de {total_pages}"
    p.font.size = Pt(9)
    p.font.name = font_name
    p.font.color.rgb = hex_to_rgb(text_hex)

def add_halftone_pattern(slide, start_x, start_y, cols=6, rows=5, dot_size=0.08, spacing=0.25, dot_color="94A3B8"):
    """Renderiza un patrón decorativo de trama de puntos (halftone) estilo diseño Canva/Impreso."""
    for r in range(rows):
        for c in range(cols):
            x = start_x + c * Inches(spacing)
            y = start_y + r * Inches(spacing)
            dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(dot_size), Inches(dot_size))
            dot.fill.solid()
            dot.fill.fore_color.rgb = hex_to_rgb(dot_color)
            dot.line.fill.background()

def add_sticky_note(slide, left, top, width, height, title, body, bg_hex="FEF08A", accent_hex="D97706"):
    """Renderiza una nota adhesiva (Post-It) estilo nota escrita a mano."""
    note = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    note.fill.solid()
    note.fill.fore_color.rgb = hex_to_rgb(bg_hex)
    note.line.color.rgb = hex_to_rgb(accent_hex)
    note.line.width = Pt(1.5)
    add_drop_shadow(note, offset_x=3, offset_y=3, blur=6, opacity=0.2, color_hex="78716C")

    # Cinta adhesiva transparente en la parte superior
    tape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + width/2 - Inches(0.6), top - Inches(0.12), Inches(1.2), Inches(0.25))
    tape.fill.solid()
    tape.fill.fore_color.rgb = hex_to_rgb("FDE047")
    tape.line.fill.background()

    txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.2), width - Inches(0.3), height - Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"📌 {title}"
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.name = "Segoe Print"
    p.font.color.rgb = hex_to_rgb(accent_hex)

    if body:
        p2 = tf.add_paragraph()
        p2.text = body
        p2.font.size = Pt(13)
        p2.font.italic = True
        p2.font.name = "Georgia"
        p2.font.color.rgb = hex_to_rgb("44403C")
        p2.space_before = Pt(8)
    return note

def add_card(slide, left, top, width, height, bg_hex, border_hex=None, border_width=1.5, shadow=True):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = hex_to_rgb(bg_hex)
    if border_hex:
        card.line.color.rgb = hex_to_rgb(border_hex)
        card.line.width = Pt(border_width)
    else:
        card.line.fill.background()
    if shadow:
        add_drop_shadow(card, offset_x=2, offset_y=2, blur=4, opacity=0.12)
    return card

def add_pill(slide, left, top, width, height, text, bg_hex, text_hex, font_name="Calibri"):
    pill = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    pill.fill.solid()
    pill.fill.fore_color.rgb = hex_to_rgb(bg_hex)
    pill.line.fill.background()
    tf = pill.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = str(text).upper()
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = hex_to_rgb(text_hex)
    p.font.name = font_name
    return pill

def generar_vector_tematico(categoria="ciberseguridad", subtema="virus", output_path=None):
    """Genera una ilustración vectorial HD (800x800 px) temática realista en PIL con gráficos de alto impacto."""
    try:
        from PIL import Image, ImageDraw
        if not output_path:
            out_dir = os.path.abspath("scratch/vectores_generados")
            os.makedirs(out_dir, exist_ok=True)
            clean_cat = "".join([c for c in str(categoria) if c.isalnum()])
            output_path = os.path.join(out_dir, f"vector_{clean_cat or 'tech'}.png")
        
        cat = str(categoria).lower()
        img = Image.new("RGBA", (800, 800), (15, 23, 42, 255))
        draw = ImageDraw.Draw(img)
        
        # 1. Trama de fondo futurista
        for x in range(40, 760, 40):
            for y in range(40, 760, 40):
                draw.ellipse([x, y, x+3, y+3], fill=(51, 65, 85, 150))
                
        if any(k in cat for k in ["ciber", "virus", "seguridad", "amenaza", "hack"]):
            # --- VECTOR REAL DE CIBERSEGURIDAD / VIRUS INFORMÁTICO ---
            # Escudo protector hexagonal
            draw.polygon([(400, 100), (650, 220), (650, 520), (400, 700), (150, 520), (150, 220)], fill=(30, 41, 59), outline=(6, 182, 212), width=6)
            draw.polygon([(400, 140), (610, 240), (610, 490), (400, 650), (190, 490), (190, 240)], outline=(236, 72, 153), width=4)
            
            # Candado / Malware Central
            draw.rectangle([340, 360, 460, 480], fill=(236, 72, 153), outline=(255, 255, 255), width=3)
            draw.arc([360, 280, 440, 380], start=180, end=360, fill=(6, 182, 212), width=6)
            
            # Rayos y alertas
            draw.line([400, 200, 400, 260], fill=(250, 204, 21), width=4)
            draw.line([250, 320, 310, 350], fill=(6, 182, 212), width=3)
            draw.line([550, 320, 490, 350], fill=(6, 182, 212), width=3)
            draw.text((330, 520), "SHIELD ACTIVE", fill=(255, 255, 255))
            
        elif any(k in cat for k in ["negocio", "startup", "finanz", "mercado"]):
            # --- VECTOR REAL DE NEGOCIOS / CRECIMIENTO ---
            draw.rectangle([100, 600, 200, 700], fill=(79, 70, 229))
            draw.rectangle([250, 450, 350, 700], fill=(6, 182, 212))
            draw.rectangle([400, 300, 500, 700], fill=(16, 185, 129))
            draw.rectangle([550, 150, 650, 700], fill=(236, 72, 153))
            draw.line([150, 550, 300, 400, 450, 250, 600, 100], fill=(250, 204, 21), width=8)
            draw.polygon([(600, 80), (630, 120), (580, 130)], fill=(250, 204, 21))
            draw.text((310, 720), "GROWTH MATRIX", fill=(255, 255, 255))
            
        else:
            # --- VECTOR TECNOLÓGICO GENERAL ---
            draw.ellipse([200, 200, 600, 600], outline=(6, 182, 212), width=6)
            draw.ellipse([280, 280, 520, 520], fill=(30, 41, 59), outline=(236, 72, 153), width=4)
            draw.text((340, 390), "ATLAS VECTOR CORE", fill=(255, 255, 255))
            
        img.save(output_path)
        return output_path
    except Exception:
        return None

def add_image_card(slide, left, top, width, height, image_path=None, label="[🖼️ RECURSO VECTORIAL / MARCA]", border_hex="38BDF8"):
    """Inserta una imagen/vector real si existe, o genera un vector temático dinámico automáticamente."""
    target_path = image_path
    if not target_path or not os.path.exists(target_path):
        cat_key = label if label else "ciberseguridad"
        target_path = generar_vector_tematico(categoria=cat_key)

    if target_path and os.path.exists(target_path):
        try:
            pic = slide.shapes.add_picture(target_path, left, top, width=width, height=height)
            return pic
        except Exception:
            pass

    card = add_card(slide, left, top, width, height, "1E293B", border_hex=border_hex, border_width=2.0)
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + height/2 - Inches(0.6), width - Inches(0.4), Inches(1.2))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "🖼️"
    p.alignment = PP_ALIGN.CENTER
    p.font.size = Pt(36)
    return card

def build_pptx_from_json(json_data):
    """
    Construye una presentación PowerPoint (.pptx) declarativa a partir de un objeto JSON especificado.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333) # 16:9 Panorámico
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    nombre_tema = json_data.get("tema", "dark_bento").lower().replace(" ", "_")
    tema = TEMAS.get(nombre_tema, TEMAS["dark_bento"])

    # Sobrescribir tema con colores personalizados si se proporcionan
    if "custom_colors" in json_data:
        custom = json_data["custom_colors"]
        tema["bg"] = custom.get("bg", tema["bg"])
        tema["card_bg"] = custom.get("card_bg", tema["card_bg"])
        tema["card_border"] = custom.get("card_border", tema["card_border"])
        tema["text_primary"] = custom.get("text_primary", tema["text_primary"])
        tema["accent_1"] = custom.get("accent_1", tema["accent_1"])

    font_family = json_data.get("font_family", tema["font_family"])
    transicion_global = json_data.get("transicion") or json_data.get("metadata", {}).get("transicion") or "fade"

    diapositivas = json_data.get("diapositivas", [])
    total_slides = len(diapositivas)
    titulo_pres = json_data.get("metadata", {}).get("titulo", "ATLAS Presentation")

    for idx, item in enumerate(diapositivas):
        tipo = item.get("tipo", "portada").lower()
        trans_type = item.get("transicion", transicion_global)
        slide = prs.slides.add_slide(blank_layout)
        add_solid_bg(slide, prs, tema["bg"], slide_idx=idx, accent_hex=tema["card_border"], accent_sec_hex=tema["accent_1"])
        apply_transition(slide, trans_type)
        add_footer(slide, prs, idx + 1, total_slides, titulo_pres, font_family, tema["text_muted"])

        # ----------------------------------------------------
        # TIPO 1: PORTADA HERO CON TARJETA MACRO
        # ----------------------------------------------------
        if tipo == "portada":
            tag = item.get("tag", "PRESENTACIÓN EJECUTIVA").upper()
            titulo = item.get("titulo", "Título Principal de Impacto")
            subtitulo = item.get("subtitulo", "Subtítulo o Resumen Ejecutivo del Proyecto")
            img_path = item.get("imagen_path")

            # 1. Header de Categoría Superior
            tx_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.4))
            p_cat = tx_cat.text_frame.paragraphs[0]
            p_cat.text = f"01. {tag}"
            p_cat.font.size = Pt(11)
            p_cat.font.bold = True
            p_cat.font.name = font_family
            p_cat.font.color.rgb = hex_to_rgb(tema["accent_1"])

            if img_path and os.path.exists(img_path):
                # Layout Split 60/40 con Imagen Real / Vector
                card_hero = add_card(slide, Inches(0.8), Inches(1.2), Inches(7.5), Inches(5.5), tema["card_bg"], tema["accent_1"], border_width=2.0)
                tx_hero = slide.shapes.add_textbox(Inches(1.1), Inches(1.6), Inches(6.9), Inches(4.7))
                tf_h = tx_hero.text_frame
                tf_h.word_wrap = True
                ph0 = tf_h.paragraphs[0]
                ph0.text = titulo
                ph0.font.size = Pt(32)
                ph0.font.bold = True
                ph0.font.name = font_family
                ph0.font.color.rgb = hex_to_rgb(tema["text_primary"])

                if subtitulo:
                    ph1 = tf_h.add_paragraph()
                    ph1.text = subtitulo
                    ph1.font.size = Pt(14)
                    ph1.font.name = font_family
                    ph1.font.color.rgb = hex_to_rgb(tema["text_muted"])
                    ph1.space_before = Pt(16)

                add_image_card(slide, Inches(8.6), Inches(1.2), Inches(3.933), Inches(5.5), image_path=img_path, border_hex=tema["accent_2"])
            else:
                # Layout Hero Dominante Ancho Completo
                card_hero = add_card(slide, Inches(0.8), Inches(1.2), Inches(11.733), Inches(5.5), tema["card_bg"], tema["accent_1"], border_width=2.0)
                tx_hero = slide.shapes.add_textbox(Inches(1.3), Inches(1.8), Inches(10.733), Inches(4.3))
                tf_h = tx_hero.text_frame
                tf_h.word_wrap = True
                ph0 = tf_h.paragraphs[0]
                ph0.text = titulo
                ph0.font.size = Pt(36)
                ph0.font.bold = True
                ph0.font.name = font_family
                ph0.font.color.rgb = hex_to_rgb(tema["text_primary"])

                if subtitulo:
                    ph1 = tf_h.add_paragraph()
                    ph1.text = subtitulo
                    ph1.font.size = Pt(16)
                    ph1.font.name = font_family
                    ph1.font.color.rgb = hex_to_rgb(tema["text_muted"])
                    ph1.space_before = Pt(16)

        # ----------------------------------------------------
        # TIPO 2: BENTO GRID CON BADGES NUMERADOS (SISTEMA DE RETÍCULA)
        # ----------------------------------------------------
        elif tipo in ["bento_grid", "tarjetas"]:
            tag = item.get("tag", "ESTRUCTURA Y COMPONENTES").upper()
            titulo = item.get("titulo", "Puntos Clave del Proyecto")
            tarjetas = item.get("tarjetas", [])

            # Header & Title 2-Level
            tx_cat = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(0.4))
            p_c = tx_cat.text_frame.paragraphs[0]
            p_c.text = f"{idx+1:02d}. {tag}"
            p_c.font.size = Pt(11)
            p_c.font.bold = True
            p_c.font.name = font_family
            p_c.font.color.rgb = hex_to_rgb(tema["accent_1"])

            tx_title = slide.shapes.add_textbox(Inches(0.8), Inches(0.85), Inches(11.733), Inches(0.8))
            p_t = tx_title.text_frame.paragraphs[0]
            p_t.text = titulo
            p_t.font.size = Pt(26)
            p_t.font.bold = True
            p_t.font.name = font_family
            p_t.font.color.rgb = hex_to_rgb(tema["text_primary"])

            num_cards = min(len(tarjetas), 3)
            if num_cards > 0:
                card_w = Inches(11.733 / num_cards - 0.3)
                spacing = Inches(0.4)
                start_x = Inches(0.8)

                for i in range(num_cards):
                    t = tarjetas[i]
                    x = start_x + i * (card_w + spacing)
                    color_acc = t.get("color", tema[f"accent_{(i%3)+1}"])

                    # Base Card
                    card = add_card(slide, x, Inches(2.2), card_w, Inches(4.5), tema["card_bg"], color_acc, border_width=1.5)

                    # Inner Badge Centrado
                    badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + card_w/2 - Inches(0.6), Inches(2.5), Inches(1.2), Inches(1.0))
                    badge.fill.solid()
                    badge.fill.fore_color.rgb = hex_to_rgb(tema["card_border"])
                    badge.line.fill.background()
                    
                    p_num = badge.text_frame.paragraphs[0]
                    p_num.text = f"0{i+1}"
                    p_num.alignment = PP_ALIGN.CENTER
                    p_num.font.size = Pt(32)
                    p_num.font.bold = True
                    p_num.font.name = font_family
                    p_num.font.color.rgb = hex_to_rgb(color_acc)

                    # Text Box con Padding seguro (Izquierda para legibilidad)
                    tb = slide.shapes.add_textbox(x + Inches(0.25), Inches(3.7), card_w - Inches(0.5), Inches(2.8))
                    tf = tb.text_frame
                    tf.word_wrap = True

                    p_t = tf.paragraphs[0]
                    p_t.text = t.get("titulo", f"Componente {i+1}")
                    p_t.alignment = PP_ALIGN.LEFT
                    p_t.font.size = Pt(16)
                    p_t.font.bold = True
                    p_t.font.name = font_family
                    p_t.font.color.rgb = hex_to_rgb(tema["text_primary"])

                    for p_item in t.get("puntos", [t.get("descripcion", "")]):
                        if p_item:
                            pi = tf.add_paragraph()
                            pi.text = f"•  {p_item}"
                            pi.alignment = PP_ALIGN.LEFT
                            pi.font.size = Pt(12)
                            pi.font.name = font_family
                            pi.font.color.rgb = hex_to_rgb(tema["text_muted"])
                            pi.space_before = Pt(8)
            tag = item.get("tag", "GRID DE CONTENIDO").upper()
            titulo = item.get("titulo", "Estructura Principal")
            tarjetas = item.get("tarjetas", [])

            add_pill(slide, Inches(1.0), Inches(0.8), Inches(2.8), Inches(0.35), tag, tema["accent_2"], tema["bg"], font_family)
            txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(11.333), Inches(0.8))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = titulo
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.name = font_family
            p.font.color.rgb = hex_to_rgb(tema["text_primary"])

            num_cards = min(len(tarjetas), 4)
            if num_cards > 0:
                card_w = Inches(11.333 / num_cards - 0.3)
                spacing = Inches(0.3)
                start_x = Inches(1.0)

                for i in range(num_cards):
                    t = tarjetas[i]
                    x = start_x + i * (card_w + spacing)
                    t_title = t.get("titulo", f"Pilar {i+1}")
                    t_desc = t.get("descripcion", "")
                    accent_color = t.get("color", tema[f"accent_{(i%3)+1}"])

                    add_card(slide, x, Inches(2.4), card_w, Inches(4.3), tema["card_bg"], accent_color)

                    txBox = slide.shapes.add_textbox(x + Inches(0.25), Inches(2.7), card_w - Inches(0.5), Inches(3.7))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    
                    p = tf.paragraphs[0]
                    p.text = t_title
                    p.font.size = Pt(20)
                    p.font.bold = True
                    p.font.name = font_family
                    p.font.color.rgb = hex_to_rgb(accent_color)

                    if "puntos" in t and isinstance(t["puntos"], list):
                        for p_idx, pt_text in enumerate(t["puntos"]):
                            p_pt = tf.add_paragraph()
                            p_pt.text = f"• {pt_text}"
                            p_pt.font.size = Pt(13)
                            p_pt.font.name = font_family
                            p_pt.font.color.rgb = hex_to_rgb(tema["text_muted"])
                            p_pt.space_before = Pt(6)
                    elif t_desc:
                        p2 = tf.add_paragraph()
                        p2.text = t_desc
                        p2.font.size = Pt(14)
                        p2.font.name = font_family
                        p2.font.color.rgb = hex_to_rgb(tema["text_muted"])
                        p2.space_before = Pt(12)

        # ----------------------------------------------------
        # TIPO 3: METRICAS / KPIS
        # ----------------------------------------------------
        elif tipo in ["kpis", "metricas"]:
            tag = item.get("tag", "MÉTRICAS CLAVE").upper()
            titulo = item.get("titulo", "Resultados Destacados")
            metricas = item.get("metricas", [])

            add_pill(slide, Inches(1.0), Inches(0.8), Inches(2.8), Inches(0.35), tag, tema["accent_3"], tema["bg"], font_family)
            txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(11.333), Inches(0.8))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = titulo
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.name = font_family
            p.font.color.rgb = hex_to_rgb(tema["text_primary"])

            num_m = min(len(metricas), 3)
            if num_m > 0:
                card_w = Inches(3.6)
                spacing = Inches(0.3)
                start_x = Inches(1.0)

                for i in range(num_m):
                    m = metricas[i]
                    x = start_x + i * (card_w + spacing)
                    val = m.get("valor", "100%")
                    lbl = m.get("etiqueta", "Métrica")
                    cambio = m.get("cambio", "")
                    color = m.get("color", tema[f"accent_{(i%3)+1}"])

                    add_card(slide, x, Inches(2.4), card_w, Inches(4.2), tema["card_bg"], color, border_width=2.0)

                    txBox = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.8), card_w - Inches(0.4), Inches(3.6))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    
                    p = tf.paragraphs[0]
                    p.text = str(val)
                    p.alignment = PP_ALIGN.CENTER
                    p.font.size = Pt(48)
                    p.font.bold = True
                    p.font.name = font_family
                    p.font.color.rgb = hex_to_rgb(color)

                    if cambio:
                        p_c = tf.add_paragraph()
                        p_c.text = str(cambio)
                        p_c.alignment = PP_ALIGN.CENTER
                        p_c.font.size = Pt(14)
                        p_c.font.bold = True
                        p_c.font.name = font_family
                        p_c.font.color.rgb = hex_to_rgb(tema["accent_2"])
                        p_c.space_before = Pt(4)

                    p2 = tf.add_paragraph()
                    p2.text = str(lbl)
                    p2.alignment = PP_ALIGN.CENTER
                    p2.font.size = Pt(15)
                    p2.font.name = font_family
                    p2.font.color.rgb = hex_to_rgb(tema["text_primary"])
                    p2.space_before = Pt(12)

        # ----------------------------------------------------
        # TIPO 4: TIMELINE / ROADMAP
        # ----------------------------------------------------
        elif tipo in ["timeline", "roadmap", "proceso"]:
            tag = item.get("tag", "LÍNEA DE TIEMPO").upper()
            titulo = item.get("titulo", "Fases del Proyecto")
            pasos = item.get("pasos", [])

            add_pill(slide, Inches(1.0), Inches(0.8), Inches(2.8), Inches(0.35), tag, tema["accent_1"], tema["bg"], font_family)
            txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(11.333), Inches(0.8))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = titulo
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.name = font_family
            p.font.color.rgb = hex_to_rgb(tema["text_primary"])

            # Conector horizontal
            conn = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(3.8), Inches(10.3), Inches(0.08))
            conn.fill.solid()
            conn.fill.fore_color.rgb = hex_to_rgb(tema["card_border"])
            conn.line.fill.background()

            num_p = min(len(pasos), 4)
            if num_p > 0:
                card_w = Inches(11.333 / num_p - 0.3)
                spacing = Inches(0.3)
                start_x = Inches(1.0)

                for i in range(num_p):
                    paso = pasos[i]
                    x = start_x + i * (card_w + spacing)
                    num_str = paso.get("numero", f"0{i+1}")
                    p_title = paso.get("titulo", f"Fase {i+1}")
                    p_desc = paso.get("descripcion", "")
                    color = paso.get("color", tema[f"accent_{(i%3)+1}"])

                    add_card(slide, x, Inches(2.4), card_w, Inches(4.3), color)

                    txBox = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.6), card_w - Inches(0.4), Inches(3.8))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    
                    p = tf.paragraphs[0]
                    p.text = num_str
                    p.font.size = Pt(36)
                    p.font.bold = True
                    p.font.name = font_family
                    p.font.color.rgb = hex_to_rgb("FFFFFF")

                    p2 = tf.add_paragraph()
                    p2.text = p_title
                    p2.font.size = Pt(18)
                    p2.font.bold = True
                    p2.font.name = font_family
                    p2.font.color.rgb = hex_to_rgb("FFFFFF")
                    p2.space_before = Pt(6)

                    p3 = tf.add_paragraph()
                    p3.text = p_desc
                    p3.font.size = Pt(13)
                    p3.font.name = font_family
                    p3.font.color.rgb = hex_to_rgb("F8FAFC")
                    p3.space_before = Pt(6)

        # ----------------------------------------------------
        # TIPO 5: VS COMPARATIVA / 2 COLUMNAS ENFRENTADAS
        # ----------------------------------------------------
        elif tipo in ["vs_comparativa", "comparativa", "foda", "pros_cons"]:
            tag = item.get("tag", "COMPARATIVA").upper()
            titulo = item.get("titulo", "Análisis Comparativo")
            col_izq = item.get("columna_izq", {})
            col_der = item.get("columna_der", {})

            add_pill(slide, Inches(1.0), Inches(0.8), Inches(2.8), Inches(0.35), tag, tema["accent_1"], tema["bg"], font_family)
            txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(11.333), Inches(0.8))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = titulo
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.name = font_family
            p.font.color.rgb = hex_to_rgb(tema["text_primary"])

            # 2 Tarjetas anchas a la izquierda y derecha
            w_col = Inches(5.4)
            add_card(slide, Inches(1.0), Inches(2.3), w_col, Inches(4.5), tema["card_bg"], tema["accent_1"], border_width=2.0)
            add_card(slide, Inches(6.933), Inches(2.3), w_col, Inches(4.5), tema["card_bg"], tema["accent_2"], border_width=2.0)

            # Columna 1 (Izquierda)
            tx_izq = slide.shapes.add_textbox(Inches(1.2), Inches(2.5), Inches(5.0), Inches(4.1))
            tf_i = tx_izq.text_frame
            tf_i.word_wrap = True
            p_i = tf_i.paragraphs[0]
            p_i.text = col_izq.get("titulo", "Opción A")
            p_i.font.size = Pt(22)
            p_i.font.bold = True
            p_i.font.name = font_family
            p_i.font.color.rgb = hex_to_rgb(tema["accent_1"])

            for item_text in col_izq.get("puntos", [col_izq.get("descripcion", "")]):
                if item_text:
                    pi = tf_i.add_paragraph()
                    pi.text = f"✔  {item_text}"
                    pi.font.size = Pt(14)
                    pi.font.name = font_family
                    pi.font.color.rgb = hex_to_rgb(tema["text_muted"])
                    pi.space_before = Pt(10)

            # Columna 2 (Derecha)
            tx_der = slide.shapes.add_textbox(Inches(7.133), Inches(2.5), Inches(5.0), Inches(4.1))
            tf_d = tx_der.text_frame
            tf_d.word_wrap = True
            p_d = tf_d.paragraphs[0]
            p_d.text = col_der.get("titulo", "Opción B")
            p_d.font.size = Pt(22)
            p_d.font.bold = True
            p_d.font.name = font_family
            p_d.font.color.rgb = hex_to_rgb(tema["accent_2"])

            for item_text in col_der.get("puntos", [col_der.get("descripcion", "")]):
                if item_text:
                    pd = tf_d.add_paragraph()
                    pd.text = f"★  {item_text}"
                    pd.font.size = Pt(14)
                    pd.font.name = font_family
                    pd.font.color.rgb = hex_to_rgb(tema["text_muted"])
                    pd.space_before = Pt(10)

        # ----------------------------------------------------
        # TIPO 6: LISTA DE FILAS / PASOS HORIZONTALES APILADOS
        # ----------------------------------------------------
        elif tipo in ["lista_pasos", "lista_filas", "puntos_clave"]:
            tag = item.get("tag", "PUNTOS CLAVE").upper()
            titulo = item.get("titulo", "Aspectos Destacados")
            filas = item.get("filas", [])

            add_pill(slide, Inches(1.0), Inches(0.8), Inches(2.8), Inches(0.35), tag, tema["accent_2"], tema["bg"], font_family)
            txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(11.333), Inches(0.8))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = titulo
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.name = font_family
            p.font.color.rgb = hex_to_rgb(tema["text_primary"])

            num_f = min(len(filas), 4)
            if num_f > 0:
                row_h = Inches(4.3 / num_f - 0.15)
                start_y = Inches(2.3)

                for i in range(num_f):
                    f_item = filas[i]
                    y = start_y + i * (row_h + Inches(0.15))
                    f_title = f_item.get("titulo", f"Punto {i+1}")
                    f_desc = f_item.get("descripcion", "")
                    color = f_item.get("color", tema[f"accent_{(i%3)+1}"])

                    add_card(slide, Inches(1.0), y, Inches(11.333), row_h, tema["card_bg"], color, border_width=1.5)

                    # Pill numérico lateral
                    add_pill(slide, Inches(1.2), y + Inches(0.15), Inches(0.8), Inches(row_h - Inches(0.3)), f"0{i+1}", color, "FFFFFF", font_family)

                    txBox = slide.shapes.add_textbox(Inches(2.2), y + Inches(0.1), Inches(9.8), row_h - Inches(0.2))
                    tf = txBox.text_frame
                    tf.word_wrap = True
                    p = tf.paragraphs[0]
                    p.text = f_title
                    p.font.size = Pt(17)
                    p.font.bold = True
                    p.font.name = font_family
                    p.font.color.rgb = hex_to_rgb(color)

                    if f_desc:
                        p2 = tf.add_paragraph()
                        p2.text = f_desc
                        p2.font.size = Pt(13)
                        p2.font.name = font_family
                        p2.font.color.rgb = hex_to_rgb(tema["text_muted"])
                        p2.space_before = Pt(4)

        # ----------------------------------------------------
        # TIPO 7: CITA / PULL QUOTE
        # ----------------------------------------------------
        elif tipo in ["cita", "pull_quote"]:
            texto_cita = item.get("cita", "“El diseño es el embajador silencioso de tu marca.”")
            autor = item.get("autor", "— Paul Rand")

            add_card(slide, Inches(1.5), Inches(1.5), Inches(10.333), Inches(4.8), tema["card_bg"], tema["card_border"])

            txBox = slide.shapes.add_textbox(Inches(2.0), Inches(2.2), Inches(9.333), Inches(3.4))
            tf = txBox.text_frame
            tf.word_wrap = True
            
            p = tf.paragraphs[0]
            p.text = texto_cita
            p.font.size = Pt(32)
            p.font.italic = True
            p.font.name = font_family
            p.font.color.rgb = hex_to_rgb(tema["text_primary"])

            if autor:
                p2 = tf.add_paragraph()
                p2.text = autor
                p2.font.size = Pt(18)
                p2.font.bold = True
                p2.font.name = font_family
                p2.font.color.rgb = hex_to_rgb(tema["accent_1"])
                p2.space_before = Pt(24)

        # ----------------------------------------------------
        # TIPO 8: NOTA ADHESIVA / POST-IT / APUNTE MANUSCRITO
        # ----------------------------------------------------
        elif tipo in ["nota_adhesiva", "post_it", "nota_mano"]:
            tag = item.get("tag", "NOTAS DE TRABAJO").upper()
            titulo = item.get("titulo", "Apuntes y Consideraciones")
            notas = item.get("notas", [])

            add_pill(slide, Inches(1.0), Inches(0.8), Inches(2.8), Inches(0.35), tag, tema["accent_3"], tema["bg"], font_family)
            txBox = slide.shapes.add_textbox(Inches(1.0), Inches(1.3), Inches(11.333), Inches(0.8))
            tf = txBox.text_frame
            p = tf.paragraphs[0]
            p.text = titulo
            p.font.size = Pt(28)
            p.font.bold = True
            p.font.name = font_family
            p.font.color.rgb = hex_to_rgb(tema["text_primary"])

            num_n = min(len(notas), 3)
            if num_n > 0:
                note_w = Inches(3.6)
                spacing = Inches(0.3)
                start_x = Inches(1.0)
                colors_postit = ["FEF08A", "FFEDD5", "E0F2FE"]
                colors_accent = ["D97706", "EA580C", "0284C7"]

                for i in range(num_n):
                    n_item = notas[i]
                    x = start_x + i * (note_w + spacing)
                    n_title = n_item.get("titulo", f"Nota {i+1}")
                    n_body = n_item.get("descripcion", "")
                    c_bg = colors_postit[i % 3]
                    c_acc = colors_accent[i % 3]

                    add_sticky_note(slide, x, Inches(2.3), note_w, Inches(4.3), n_title, n_body, bg_hex=c_bg, accent_hex=c_acc)

        # ----------------------------------------------------
        # TIPO 8: CIERRE / CALL TO ACTION / FIN
        # ----------------------------------------------------
        elif tipo in ["cierre", "fin", "cta", "final", "cierre_cta"]:
            tag = item.get("tag", "CIERRE").upper()
            titulo = item.get("titulo", "Gracias")
            subtitulo = item.get("subtitulo", "")

            # Fondo de color sólido usando el accent_1
            bg_cierre = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
            bg_cierre.fill.solid()
            bg_cierre.fill.fore_color.rgb = hex_to_rgb(tema["accent_1"])
            bg_cierre.line.fill.background()

            # Franjas decorativas superior e inferior
            franja_top = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.15))
            franja_top.fill.solid()
            franja_top.fill.fore_color.rgb = hex_to_rgb("FFFFFF")
            franja_top.line.fill.background()

            franja_bot = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.15),
                prs.slide_width, Inches(0.15))
            franja_bot.fill.solid()
            franja_bot.fill.fore_color.rgb = hex_to_rgb("FFFFFF")
            franja_bot.line.fill.background()

            # Pill centrado
            add_pill(slide, Inches(4.9), Inches(1.4), Inches(3.5), Inches(0.42),
                     tag, tema["card_bg"], tema["accent_1"], font_family)

            # Título grande centrado en blanco
            txBox = slide.shapes.add_textbox(
                Inches(1.0), Inches(2.1), Inches(11.333), Inches(2.4))
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = titulo
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(48)
            p.font.bold = True
            p.font.name = font_family
            p.font.color.rgb = hex_to_rgb("FFFFFF")

            # Línea decorativa horizontal
            sep = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(5.1), Inches(4.6), Inches(3.133), Inches(0.05))
            sep.fill.solid()
            sep.fill.fore_color.rgb = hex_to_rgb("FFFFFF")
            sep.line.fill.background()

            # Subtítulo / contacto centrado
            if subtitulo:
                txBox2 = slide.shapes.add_textbox(
                    Inches(1.5), Inches(4.8), Inches(10.333), Inches(1.5))
                tf2 = txBox2.text_frame
                tf2.word_wrap = True
                p3 = tf2.paragraphs[0]
                p3.text = subtitulo
                p3.alignment = PP_ALIGN.CENTER
                p3.font.size = Pt(16)
                p3.font.name = font_family
                p3.font.color.rgb = hex_to_rgb("FFFFFF")

    return prs

def parse_markdown_to_pptx_json(md_text, tema="dark_bento"):
    """
    Parsea un texto en Markdown o borrador simple y lo convierte automáticamente
    en una estructura JSON equivalente para build_pptx_from_json.
    """
    lines = [l.strip() for l in md_text.splitlines() if l.strip()]
    diapositivas = []
    
    current_slide = None
    title_main = "Presentación Creada por ATLAS"
    subtitle_main = ""

    for line in lines:
        if line.startswith("# "):
            title_main = line.lstrip("# ").strip()
        elif line.startswith("## "):
            if current_slide:
                diapositivas.append(current_slide)
            slide_title = line.lstrip("## ").strip()
            current_slide = {
                "tipo": "bento_grid",
                "tag": "SECCIÓN",
                "titulo": slide_title,
                "tarjetas": []
            }
        elif (line.startswith("- ") or line.startswith("* ")) and current_slide:
            item_text = line[2:].strip()
            parts = item_text.split(":", 1)
            if len(parts) == 2:
                t_head, t_desc = parts[0].strip(), parts[1].strip()
            else:
                t_head, t_desc = "Punto Clave", parts[0].strip()
            
            current_slide["tarjetas"].append({
                "titulo": t_head,
                "descripcion": t_desc
            })

    if current_slide:
        diapositivas.append(current_slide)

    # Insertar Portada al inicio
    diapositivas.insert(0, {
        "tipo": "portada",
        "tag": "PRESENTACIÓN AUTOMÁTICA",
        "titulo": title_main,
        "subtitulo": subtitle_main or "Generada automáticamente desde documento de texto"
    })

    return {
        "tema": tema,
        "diapositivas": diapositivas
    }

def compile_yaml_to_pptx(yaml_path: str, pptx_path: str = None) -> str:
    """
    Función pública estandarizada para compilar una especificación YAML, JSON o MD
    a un archivo de presentación PowerPoint (.pptx).
    """
    input_path = os.path.abspath(yaml_path)
    if not os.path.exists(input_path):
        raise FileNotFoundError(f"Archivo de especificación no encontrado: {input_path}")

    if not pptx_path:
        pptx_path = os.path.splitext(input_path)[0] + ".pptx"
    output_pptx = os.path.abspath(pptx_path)

    out_dir = os.path.dirname(output_pptx)
    if out_dir and not os.path.exists(out_dir):
        os.makedirs(out_dir, exist_ok=True)

    ext = os.path.splitext(input_path)[1].lower()
    if ext == ".json":
        with open(input_path, "r", encoding="utf-8") as f:
            data = json.load(f)
    elif ext in [".yaml", ".yml"]:
        import yaml
        with open(input_path, "r", encoding="utf-8") as f:
            data = yaml.safe_load(f)
    else:
        with open(input_path, "r", encoding="utf-8") as f:
            raw_text = f.read()
        data = parse_markdown_to_pptx_json(raw_text)

    # Desempaquetar la estructura si viene bajo el nodo raíz 'presentacion'
    if isinstance(data, dict) and "presentacion" in data:
        pres = data["presentacion"]
        meta = pres.get("metadata", {})
        estilos = pres.get("estilos", {})
        slides = pres.get("diapositivas", [])
        
        json_data = {
            "tema": meta.get("tema_visual", "dark_bento"),
            "font_family": estilos.get("font_family", "Calibri"),
            "diapositivas": slides
        }
        if "palette" in estilos:
            pal = estilos["palette"]
            json_data["custom_colors"] = {
                "bg": pal.get("primario"),
                "card_bg": pal.get("secundario"),
                "accent_1": pal.get("acento")
            }
    else:
        json_data = data

    prs = build_pptx_from_json(json_data)
    prs.save(output_pptx)
    return output_pptx

def main():
    parser = argparse.ArgumentParser(
        description="Generador autónomo de presentaciones PowerPoint (.pptx) de alto impacto visual (Soporta YAML, JSON, Markdown y Texto)."
    )
    parser.add_argument("input_file", help="Ruta al archivo (.yaml, .json, .md, .txt) con el contenido de la presentación")
    parser.add_argument("output_pptx", nargs="?", help="Ruta de salida para el archivo .pptx (opcional)")
    parser.add_argument("--tema", default="dark_bento", help="Tema visual (dark_bento, canva_creative, corporate_executive, neon_cyberpunk, minimal_editorial, vibrant_roadmap)")
    args = parser.parse_args()

    try:
        output = compile_yaml_to_pptx(args.input_file, args.output_pptx)
        size = os.path.getsize(output)
        print(f"✅ [OK] Presentación PowerPoint generada exitosamente: {output} ({size / 1024:.1f} KB)")
        sys.exit(0)
    except Exception as e:
        import traceback
        print(f"❌ [ERROR] Falló la compilación de la presentación: {e}", file=sys.stderr)
        traceback.print_exc()
        sys.exit(1)

if __name__ == "__main__":
    main()


